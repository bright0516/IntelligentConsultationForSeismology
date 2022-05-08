'''
*******************************
简介：自动生成会商PPT程序
作者：王光明  日期：2021/12/23
版本：test01
*******************************
'''

from pathlib import Path
from datetime import datetime
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE



def firstPage(pf):
    # 第一页PPT，主要包括题目、作者和时间
    # 设置PPT长宽
    # PPT标准长宽（4:3 宽25.4，高19.05； 16:9 宽33.867，高19.05 cm）
    sl_w = Cm(25.4)
    sl_h = Cm(19.05)
    pf.slide_width = sl_w
    pf.slide_height = sl_h
    # 新增PPT页
    # 0表示默认模板第1页，6表示空白页
    # --------------------PPT首页模板----------------------
    slide = pf.slides.add_slide(pf.slide_layouts[6])
    # 插入ppt标题文本框
    # left = top = width = height = Cm(3)
    left = Cm(0)
    top = Cm(4)
    width = Cm(25.4)
    height = Cm(4)
    text_box = slide.shapes.add_textbox(left, top, width, height)
    tf = text_box.text_frame
    p_title = tf.add_paragraph()
    p_title.text = "震 情 监 视 报 告"
    p_title.alignment = PP_ALIGN.CENTER
    p_title.vertical_anchor = MSO_ANCHOR.MIDDLE
    # 修改字体(SlidePlaceholder不能修改字体)
    p_title.font.bold = True
    p_title.font.name = "黑体"
    p_title.font.color.rgb = RGBColor(255, 255, 255)
    p_title.font.size = Pt(48)
    # 文本框边框样式调整
    fill = text_box.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(51, 102, 205)
    # line = text_box.line
    # line.color.rgb = RGBColor(0, 0, 0)
    # line.width = Cm(0.1)

    # 插入作者文本框
    # left = top = width = height = Cm(3)
    left = Cm(3)
    top = Cm(10)
    width = Cm(19.4)
    height = Cm(3)
    text_box = slide.shapes.add_textbox(left, top, width, height)
    tf = text_box.text_frame
    p_author = tf.add_paragraph()
    p_author.text = "地震学研究室"
    p_author.alignment = PP_ALIGN.CENTER
    # 修改字体(SlidePlaceholder不能修改字体)
    # p_author.font.bold = True
    p_author.font.name = "楷体"
    p_author.font.color.rgb = RGBColor(0, 0, 128)
    p_author.font.size = Pt(32)

    # 插入日期文本框
    # left = top = width = height = Cm(3)
    left = Cm(3)
    top = Cm(13)
    width = Cm(19.4)
    height = Cm(3)
    text_box = slide.shapes.add_textbox(left, top, width, height)
    tf = text_box.text_frame
    p_time = tf.add_paragraph()
    p_time.text = str(datetime.now().date())  # 使用datetime库确定当前时间
    p_time.alignment = PP_ALIGN.CENTER
    # 修改字体(SlidePlaceholder不能修改字体)
    # p_time.font.bold = True
    p_time.font.name = "Times New Roman"
    # p.font.color.rgb = RGBColor(247, 150, 70)
    p_time.font.color.rgb = RGBColor(0, 0, 128)
    p_time.font.size = Pt(28)

    return pf
# --------------------------------------------------

def secondPage(pf):
    # 新建一页空白页
    slide = pf.slides.add_slide(pf.slide_layouts[6])
    # 插入文本框
    left = Cm(0)
    top = Cm(3)
    width = Cm(25.4)
    height = Cm(3)
    text_box = slide.shapes.add_textbox(left, top, width, height)
    # 插入内容——标题
    tf = text_box.text_frame
    p_author = tf.add_paragraph()
    p_author.text = "主要内容"
    p_author.alignment = PP_ALIGN.CENTER
    # 修改字体(SlidePlaceholder不能修改字体)
    # p_author.font.bold = True
    p_author.font.name = "楷体"
    p_author.font.color.rgb = RGBColor(0, 0, 128)
    p_author.font.size = Pt(36)
    # 修改背景颜色
    # fill = text_box.fill
    # fill.solid()
    # fill.fore_color.rgb = RGBColor(51, 102, 205)

    # 插入线段-有bug
    # left = Cm(1)
    # top = Cm(6)
    # width = Cm(23.4)
    # height = Cm(0.2)
    # line = slide.shapes.add_shape(MSO_SHAPE.LINE_CALLOUT_1, left, top, width, height)
    # li = line.line
    # li.color.rgb = RGBColor(255, 0, 0)


    # 插入文本框
    left = Cm(3)
    top = Cm(7)
    width = Cm(19.4)
    height = Cm(10)
    text_box = slide.shapes.add_textbox(left, top, width, height)
    # 插入内容——目录
    tf = text_box.text_frame
    p1 = tf.add_paragraph()
    p1.text = "一、地震活动概况"
    # p_author.alignment = PP_ALIGN.CENTER
    # 修改字体(SlidePlaceholder不能修改字体)
    # p_author.font.bold = True
    p1.font.name = "宋体"
    p1.font.color.rgb = RGBColor(255, 0, 0)
    p1.font.size = Pt(28)
    p1.line_spacing = 1.5  # 1.5 倍的行距
    # -----------------------------------------
    p1 = tf.add_paragraph()
    p1.text = "二、显著地震（现象）及地震序列跟踪"
    # p_author.alignment = PP_ALIGN.CENTER
    # 修改字体(SlidePlaceholder不能修改字体)
    # p_author.font.bold = True
    p1.font.name = "宋体"
    p1.font.color.rgb = RGBColor(0, 0, 0)
    p1.font.size = Pt(28)
    p1.line_spacing = 1.5  # 1.5 倍的行距
    # -----------------------------------------
    p1 = tf.add_paragraph()
    p1.text = "三、异常跟踪分析"
    # p_author.alignment = PP_ALIGN.CENTER
    # 修改字体(SlidePlaceholder不能修改字体)
    # p_author.font.bold = True
    p1.font.name = "宋体"
    p1.font.color.rgb = RGBColor(0, 0, 0)
    p1.font.size = Pt(28)
    p1.line_spacing = 1.5  # 1.5 倍的行距
    # -----------------------------------------
    p1 = tf.add_paragraph()
    p1.text = "四、综合分析和结论"
    # p_author.alignment = PP_ALIGN.CENTER
    # 修改字体(SlidePlaceholder不能修改字体)
    # p_author.font.bold = True
    p1.font.name = "宋体"
    p1.font.color.rgb = RGBColor(0, 0, 0)
    p1.font.size = Pt(28)
    p1.line_spacing = 1.5  # 1.5 倍的行距
    return pf

# 查看ppt占位符id
# for shape in slide.placeholders:
#     phf = shape.placeholder_format
#     print(f'{phf.idx}--{shape.name}--{phf.type}')
# pf_title = slide.placeholders[0]
# pf_subtitle = slide.placeholders[1]
# pf_title.text = "地震学研究室周会商报告"
# pf_time = time.strftime("%Y-%m-%d", time.localtime(time.time()))   # 使用time库确定当前时间
# pf_subtitle.text = pf_time
# pf_subtitle.font.name = "Times New Roman"
# 设置背景颜色
# slide.background.fill.solid()
# slide.background.fill.fore_color.rgb = RGBColor(230, 230, 230)
# 保存PPT文件
def main():
    src_folder = Path('D:\\python_wgm\\ppt_for_work\\')  # 文件夹绝对路径
    pf = Presentation()  # 不带文件名，表示新建PPT文件
    # PPT标准长宽（4:3 宽25.4，高19.05； 16:9 宽33.867，高19.05 cm）
    sl_w = Cm(25.4)
    sl_h = Cm(19.05)
    pf.slide_width = sl_w
    pf.slide_height = sl_h
    # 第一页PPT
    firstPage(pf)
    # 第二页PPT
    secondPage(pf)

    pf.save('D:\\python_wgm\\ppt_for_work\\test1.pptx')

if __name__ == '__main__':
    main()
